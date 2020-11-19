import sys
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt
import sys
import os
import glob
import win32com.client

TARGET_PPT = r"C:\Users\37197\Desktop\Ineuron\Internship\GluonCV\my_first_ppt.pptx"
keyword='John Doe'
replacement='Jerryl Davis'
prs = Presentation(TARGET_PPT)
text_runs = []


def convert(files, formatType=32):
    powerpoint = win32com.client.Dispatch("Powerpoint.Application")
    powerpoint.Visible = 1
    for filename in files:
        newname = os.path.splitext(filename)[0] + ".pdf"
        print(filename)
        deck = powerpoint.Presentations.Open(filename)
        deck.SaveAs(newname, formatType)
        deck.Close()
    powerpoint.Quit()

for slide in prs.slides:
    top = width = height = Inches(3.6)
    left = Inches(5)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = str(sys.argv[1])+" "+str(sys.argv[2])
    p.font.bold = True
    p.font.size = Pt(33)
    font=p.font
    font.color.rgb = RGBColor(204, 0, 0)
    name=str(sys.argv[1])+" "+str(sys.argv[2])+".pptx"
    path="C:\\Users\\37197\\Desktop\\Ineuron\\Internship\\GluonCV\\certs"+"\\"+name
    top2 = width2 = height2 = Inches(4.2)
    left = Inches(5.4)
    txBox2 = slide.shapes.add_textbox(left, top2, width2, height2)
    tf2 = txBox2.text_frame
    p2 = tf2.add_paragraph()
    p2.text = str(sys.argv[3])
    p2.font.bold = True
    p2.bullet=False
    prs.save(path)
    files = glob.glob(path)
    convert(files)